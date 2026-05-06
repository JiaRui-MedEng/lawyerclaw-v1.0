"""
文件读取工具 - 增强版（支持文档解析）

用法:
    from service.tools.file_reader_enhanced import read_any_file
    
    result = read_any_file('test.docx')
    if result['success']:
        print(result['content'])
"""
import os
import re
from pathlib import Path
from typing import Dict, Any

# 文档文件扩展名
_DOCUMENT_EXTENSIONS = {
    '.pdf': 'PDF 文档',
    '.docx': 'Word 文档',
    '.doc': 'Word 文档 (旧版)',
    '.pptx': 'PowerPoint 演示文稿',
    '.xlsx': 'Excel 工作表',
}

# 二进制文件扩展名
_BINARY_EXTENSIONS = {
    '.exe', '.dll', '.so', '.bin', '.dat',
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.ico',
    '.zip', '.rar', '.tar', '.gz', '.7z',
    '.pyc', '.pyo', '.o', '.a', '.lib',
    '.mp3', '.mp4', '.avi', '.mov', '.mkv',
}


def parse_document(file_path: str) -> Dict[str, Any]:
    """解析文档文件（docx, pdf, pptx, xlsx）"""
    ext = Path(file_path).suffix.lower()
    
    try:
        if ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # 提取表格
            tables = []
            for table in doc.tables:
                table_text = "表格:\n"
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    table_text += f"  行{i+1}: {' | '.join(cells)}\n"
                tables.append(table_text)
            
            content = "\n\n".join(paragraphs + tables)
            return {
                'success': True,
                'content': content,
                'pages': len(paragraphs),
                'file_type': 'docx'
            }
            
        elif ext == '.pdf':
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            pages = []
            total_pages = len(doc)

            for i in range(total_pages):
                page = doc[i]
                text = page.get_text()
                if text.strip():
                    pages.append(f"--- 第 {i+1} 页 ---\n{text}")

            doc.close()

            if not pages:
                return {'success': False, 'error': 'PDF 没有可提取的文本（可能是扫描版）'}

            return {
                'success': True,
                'content': "\n\n".join(pages),
                'pages': total_pages,
                'file_type': 'pdf'
            }
            
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"--- 幻灯片 {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                slides.append(slide_text)
            
            return {
                'success': True,
                'content': "\n\n".join(slides),
                'slides': len(prs.slides),
                'file_type': 'pptx'
            }
            
        elif ext == '.xlsx':
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            sheets = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = f"--- 工作表：{sheet_name} ---\n"
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_str = ' | '.join([str(c) if c is not None else '' for c in row])
                        rows.append(row_str)
                sheet_text += "\n".join(rows[:100])  # 限制 100 行
                sheets.append(sheet_text)
            
            return {
                'success': True,
                'content': "\n\n".join(sheets),
                'sheets': len(wb.sheetnames),
                'file_type': 'xlsx'
            }
            
        else:
            return {'success': False, 'error': f'不支持的文档格式：{ext}'}
            
    except ImportError as e:
        return {'success': False, 'error': f'缺少依赖库：{str(e)}'}
    except Exception as e:
        return {'success': False, 'error': f'解析失败：{str(e)}'}


def read_any_file(file_path: str, max_chars: int = 50000) -> Dict[str, Any]:
    """
    读取任意文件（文本或文档）
    
    Args:
        file_path: 文件路径
        max_chars: 最大字符数
        
    Returns:
        {
            'success': bool,
            'path': str,
            'content': str or None,
            'error': str or None,
            'file_type': str,
            'pages': int,
            'hint': str
        }
    """
    ext = Path(file_path).suffix.lower()
    
    # 1. 检查文件是否存在
    path = Path(file_path)
    if not path.exists():
        return {'success': False, 'error': f'文件不存在：{file_path}'}
    
    if not path.is_file():
        return {'success': False, 'error': f'不是文件：{file_path}'}
    
    # 2. 检查是否是文档文件
    if ext in _DOCUMENT_EXTENSIONS:
        print(f"📄 检测到文档文件：{ext}")
        result = parse_document(str(path))
        
        if result['success']:
            content = result.get('content', '')
            
            # 添加行号
            numbered = '\n'.join([f"{i+1:6d}: {line}" for i, line in enumerate(content.split('\n'))])
            
            return {
                'success': True,
                'path': str(path),
                'content': numbered[:max_chars],
                'file_type': result.get('file_type'),
                'pages': result.get('pages', result.get('slides', result.get('sheets', 0))),
                'hint': f"[文档类型：{_DOCUMENT_EXTENSIONS.get(ext)}，{result.get('pages', 0)} 页/段落]"
            }
        else:
            return result
    
    # 3. 检查是否是二进制文件
    if ext in _BINARY_EXTENSIONS:
        return {'success': False, 'error': f'无法读取二进制文件 ({ext})'}
    
    # 4. 读取文本文件
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read(max_chars)
        
        # 添加行号
        numbered = '\n'.join([f"{i+1:6d}: {line}" for i, line in enumerate(content.split('\n'))])
        
        return {
            'success': True,
            'path': str(path),
            'content': numbered,
            'file_type': 'text',
            'hint': f"[文本文件]"
        }
        
    except UnicodeDecodeError:
        return {'success': False, 'error': '文件编码错误，不是 UTF-8 编码'}
    except Exception as e:
        return {'success': False, 'error': f'读取失败：{str(e)}'}


# 测试
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        result = read_any_file(file_path)
        
        if result['success']:
            print(f"✅ 读取成功")
            print(f"类型：{result.get('file_type')}")
            print(f"页数/段落：{result.get('pages', 0)}")
            print(f"提示：{result.get('hint')}")
            print(f"\n内容预览:\n{result['content'][:500]}")
        else:
            print(f"❌ 读取失败：{result['error']}")
    else:
        print("用法：python file_reader_enhanced.py <文件路径>")
