"""
文档解析工具

支持格式:
- .docx (Word 2007+)
- .pdf (PDF 文档)
- .doc (旧版 Word)
- .pptx (PowerPoint)
- .xlsx (Excel)
"""
import os
from pathlib import Path
from typing import Dict, Any, List


# ═══════════════════════════════════════════════════════════
# DOCX 解析
# ═══════════════════════════════════════════════════════════

def parse_docx(file_path: str) -> Dict[str, Any]:
    """
    解析 Word (.docx) 文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        {
            'success': bool,
            'content': str,
            'pages': int,  # 段落数
            'words': int,  # 字数
            'error': str or None
        }
    """
    try:
        from docx import Document
        
        doc = Document(file_path)
        
        # 提取所有段落
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        # 提取表格内容
        tables_content = []
        for table in doc.tables:
            table_text = "表格:\n"
            for i, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                table_text += f"  行{i+1}: {' | '.join(row_text)}\n"
            tables_content.append(table_text)
        
        # 合并内容
        content = "\n\n".join(paragraphs)
        if tables_content:
            content += "\n\n" + "\n\n".join(tables_content)
        
        # 统计
        full_text = doc.paragraphs[0].text if doc.paragraphs else ""
        word_count = len(full_text.split())
        
        return {
            'success': True,
            'content': content,
            'pages': len(paragraphs),
            'words': word_count,
            'file_type': 'docx'
        }
        
    except ImportError:
        return {
            'success': False,
            'error': '未安装 python-docx 库，请运行：pip install python-docx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'解析失败：{str(e)}'
        }


# ═══════════════════════════════════════════════════════════
# PDF 解析
# ═══════════════════════════════════════════════════════════

def parse_pdf(file_path: str) -> Dict[str, Any]:
    """
    解析 PDF 文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        {
            'success': bool,
            'content': str,
            'pages': int,
            'error': str or None
        }
    """
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        pages_content = []
        total_pages = len(doc)

        for page_num in range(total_pages):
            page = doc[page_num]
            text = page.get_text()

            if text.strip():
                pages_content.append(f"--- 第 {page_num + 1} 页 ---\n{text}")

        doc.close()

        if not pages_content:
            return {
                'success': False,
                'error': 'PDF 文件没有可提取的文本内容（可能是扫描版）'
            }

        return {
            'success': True,
            'content': "\n\n".join(pages_content),
            'pages': total_pages,
            'file_type': 'pdf'
        }
        
    except ImportError:
        return {
            'success': False,
            'error': '未安装 PyMuPDF 库，请运行：pip install pymupdf'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'解析失败：{str(e)}'
        }


# ═══════════════════════════════════════════════════════════
# PPTX 解析
# ═══════════════════════════════════════════════════════════

def parse_pptx(file_path: str) -> Dict[str, Any]:
    """
    解析 PowerPoint (.pptx) 文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        {
            'success': bool,
            'content': str,
            'slides': int,
            'error': str or None
        }
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = f"--- 幻灯片 {i + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += f"{shape.text.strip()}\n"
            
            slides_content.append(slide_text)
        
        return {
            'success': True,
            'content': "\n\n".join(slides_content),
            'slides': len(prs.slides),
            'file_type': 'pptx'
        }
        
    except ImportError:
        return {
            'success': False,
            'error': '未安装 python-pptx 库，请运行：pip install python-pptx'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'解析失败：{str(e)}'
        }


# ═══════════════════════════════════════════════════════════
# XLSX 解析
# ═══════════════════════════════════════════════════════════

def parse_xlsx(file_path: str) -> Dict[str, Any]:
    """
    解析 Excel (.xlsx) 文件
    
    Args:
        file_path: 文件路径
        
    Returns:
        {
            'success': bool,
            'content': str,
            'sheets': int,
            'error': str or None
        }
    """
    try:
        from openpyxl import load_workbook
        
        wb = load_workbook(file_path, data_only=True)
        sheets_content = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = f"--- 工作表：{sheet_name} ---\n"
            
            # 读取所有行
            rows = []
            for row in sheet.iter_rows(values_only=True):
                # 过滤空行
                if any(cell is not None for cell in row):
                    row_text = []
                    for cell in row:
                        row_text.append(str(cell) if cell is not None else '')
                    rows.append(' | '.join(row_text))
            
            sheet_text += "\n".join(rows[:100])  # 限制每个工作表 100 行
            sheets_content.append(sheet_text)
        
        return {
            'success': True,
            'content': "\n\n".join(sheets_content),
            'sheets': len(wb.sheetnames),
            'file_type': 'xlsx'
        }
        
    except ImportError:
        return {
            'success': False,
            'error': '未安装 openpyxl 库，请运行：pip install openpyxl'
        }
    except Exception as e:
        return {
            'success': False,
            'error': f'解析失败：{str(e)}'
        }


# ═══════════════════════════════════════════════════════════
# 统一接口
# ═══════════════════════════════════════════════════════════

def parse_document(file_path: str) -> Dict[str, Any]:
    """
    根据文件扩展名自动选择解析器
    
    Args:
        file_path: 文件路径
        
    Returns:
        解析结果
    """
    ext = Path(file_path).suffix.lower()
    
    parsers = {
        '.docx': parse_docx,
        '.pdf': parse_pdf,
        '.pptx': parse_pptx,
        '.xlsx': parse_xlsx,
    }
    
    if ext in parsers:
        return parsers[ext](file_path)
    else:
        return {
            'success': False,
            'error': f'不支持的文件格式：{ext}\n支持的格式：.docx, .pdf, .pptx, .xlsx'
        }


# ═══════════════════════════════════════════════════════════
# 测试
# ═══════════════════════════════════════════════════════════

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        file_path = sys.argv[1]
        result = parse_document(file_path)
        
        if result['success']:
            print(f"✅ 解析成功")
            print(f"类型：{result.get('file_type')}")
            print(f"页数/段落数：{result.get('pages', result.get('slides', result.get('sheets', 0)))}")
            print(f"\n内容预览 (前 500 字符):\n{result['content'][:500]}")
        else:
            print(f"❌ 解析失败：{result['error']}")
    else:
        print("用法：python document_parser.py <文件路径>")
